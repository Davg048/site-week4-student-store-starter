"""Builds StudentStore-Demo.pptx — a 5-slide deck for the CodePath backend demo.
Targets the official demo rubric: show MODELS via Prisma Studio, ROUTES via Postman.
Run: python3 build_slides.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- palette ----
NAVY   = RGBColor(0x0F, 0x1B, 0x33)
INK    = RGBColor(0x1A, 0x20, 0x2C)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
MUTED  = RGBColor(0xB6, 0xC2, 0xD9)
ACCENT = RGBColor(0x3B, 0x82, 0xF6)  # blue
GREEN  = RGBColor(0x16, 0xA3, 0x4A)
YELLOW = RGBColor(0xCA, 0x8A, 0x04)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
CYAN   = RGBColor(0x06, 0xB6, 0xD4)
CODEBG = RGBColor(0x0B, 0x10, 0x21)
CODEFG = RGBColor(0xE6, 0xED, 0xF3)
GREENC = RGBColor(0x7E, 0xE7, 0x87)

SANS = "Calibri"
MONO = "Consolas"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def box(slide, l, t, w, h, fill=None, line=None, line_w=None, radius=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, l, t, w, h)
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = line_w or Pt(1)
    shp.shadow.inherit = False
    return shp


def text(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.0):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after); p.line_spacing = line_spacing
        for (s, size, color, bold, italic, font) in para:
            r = p.add_run(); r.text = s
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.italic = italic; r.font.name = font
    return tb


def chip(slide, l, t, label, color, w=2.0):
    c = box(slide, l, t, Inches(w), Inches(0.45), fill=color, radius=True)
    tf = c.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = MONO
    return c

# ============================================================ SLIDE 1 — TITLE
s = prs.slides.add_slide(BLANK); bg(s, NAVY)
box(s, 0, Inches(2.55), SW, Inches(0.06), fill=ACCENT)
text(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(1.2),
     [[("🏪  Student Store — Backend API", 46, WHITE, True, False, SANS)]])
text(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(0.6),
     [[("A REST API + relational database + atomic checkout, built from scratch",
        20, MUTED, False, True, SANS)]])
text(s, Inches(0.9), Inches(3.5), Inches(11.5), Inches(0.5),
     [[("by [Your Name]  ·  [pronouns]", 18, WHITE, False, False, SANS)]])
for i, (lbl, col) in enumerate([("Node + Express", ACCENT), ("Prisma ORM", GREEN),
                                 ("PostgreSQL", YELLOW), ("Postman-tested", PURPLE)]):
    chip(s, Inches(0.9 + i*2.4), Inches(4.4), lbl, col, w=2.15)
text(s, Inches(0.9), Inches(5.5), Inches(11.5), Inches(0.5),
     [[("Demo plan:  Models in Prisma Studio  →  Routes in Postman  →  the transaction",
        16, GREENC, False, False, MONO)]])
text(s, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.5),
     [[("CodePath SITE · Unit 4 · 4–5 min + 2 min Q&A", 13, MUTED, False, False, SANS)]])

# ============================================================ SLIDE 2 — MODELS (Prisma Studio)
s = prs.slides.add_slide(BLANK); bg(s, NAVY)
text(s, Inches(0.9), Inches(0.45), Inches(11.6), Inches(0.9),
     [[("🗂️  The Data Models  ", 34, WHITE, True, False, SANS),
       ("— show live in  npx prisma studio", 20, GREENC, False, False, MONO)]])
# three model cards
mw = Inches(3.85); mx = Inches(0.7); my = Inches(1.6); mg = Inches(0.18)
models = [
    ("Product", ACCENT, ["id  (PK)", "name", "description", "price", "imageUrl ?", "category"]),
    ("Order", PURPLE, ["id  (PK)", "customer", "totalPrice", "status  =\"pending\"", "createdAt = now()"]),
    ("OrderItem", CYAN, ["id  (PK)", "orderId  (FK→Order)", "productId (FK→Product)", "quantity", "price (snapshot)"]),
]
for i,(name,col,fields) in enumerate(models):
    bx = Emu(int(mx)+i*(int(mw)+int(mg)))
    box(s, bx, my, mw, Inches(3.3), fill=INK, line=col, line_w=Pt(2), radius=True)
    box(s, bx, my, mw, Inches(0.55), fill=col, radius=True)
    text(s, Emu(int(bx)+Emu(150000)), my, Inches(3.5), Inches(0.55),
         [[(name, 18, WHITE, True, False, MONO)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Emu(int(bx)+Emu(220000)), Emu(int(my)+Emu(680000)), Inches(3.4), Inches(2.6),
         [[("• "+f, 14, WHITE, False, False, MONO)] for f in fields], space_after=5)
text(s, Inches(0.7), Inches(5.15), Inches(11.9), Inches(1.9),
     [[("How they connect:", 18, GREEN, True, False, SANS)],
      [("OrderItem is a JOIN table — it links one Order to one Product (two foreign keys).",
        15, WHITE, False, False, SANS)],
      [("Cascade delete: removing a Product OR an Order auto-removes its OrderItems — enforced by Postgres.",
        15, WHITE, False, False, SANS)],
      [("Say: \"Each field maps to a real column. camelCase in code, snake_case in the DB, bridged with @map.\"",
        13, MUTED, False, True, SANS)]],
     space_after=6)

# ============================================================ SLIDE 3 — ROUTES (Postman)
s = prs.slides.add_slide(BLANK); bg(s, NAVY)
text(s, Inches(0.9), Inches(0.45), Inches(11.6), Inches(0.9),
     [[("🛣️  The Routes  ", 34, WHITE, True, False, SANS),
       ("— walk the Postman collection", 20, GREENC, False, False, MONO)]])
# endpoint list (two columns)
left = [
    ("GET",    "/products", ACCENT),
    ("GET",    "/products/:id", ACCENT),
    ("POST",   "/products", GREEN),
    ("PUT",    "/products/:id", YELLOW),
    ("DELETE", "/products/:id", RGBColor(0xDC,0x26,0x26)),
]
right = [
    ("GET",    "/orders", ACCENT),
    ("GET",    "/orders/:id", ACCENT),
    ("POST",   "/orders  ⭐", GREEN),
    ("PUT",    "/orders/:id", YELLOW),
    ("DELETE", "/orders/:id", RGBColor(0xDC,0x26,0x26)),
]
def verb_rows(slide, items, x, y):
    for i,(verb,path,col) in enumerate(items):
        ry = Emu(int(y)+i*int(Inches(0.62)))
        vb = box(slide, x, ry, Inches(1.25), Inches(0.5), fill=col, radius=True)
        tf=vb.text_frame; tf.vertical_anchor=MSO_ANCHOR.MIDDLE; p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
        r=p.add_run(); r.text=verb; r.font.size=Pt(12); r.font.bold=True; r.font.color.rgb=WHITE; r.font.name=MONO
        text(slide, Emu(int(x)+int(Inches(1.4))), ry, Inches(4.0), Inches(0.5),
             [[(path, 16, WHITE, False, False, MONO)]], anchor=MSO_ANCHOR.MIDDLE)
verb_rows(s, left, Inches(0.9), Inches(1.6))
verb_rows(s, right, Inches(6.9), Inches(1.6))
text(s, Inches(0.9), Inches(4.9), Inches(5.5), Inches(0.5),
     [[("+ /order-items  ·  /orders/:id/items", 13, MUTED, False, True, MONO)]])
# demo punchline box
box(s, Inches(0.7), Inches(5.5), Inches(11.9), Inches(1.5), fill=INK, line=GREEN, line_w=Pt(2), radius=True)
text(s, Inches(1.0), Inches(5.65), Inches(11.4), Inches(1.3),
     [[("⭐ The money-shot in Postman:", 16, GREENC, True, False, SANS)],
      [("Send POST /orders with a fake product → 404. Then GET /orders → nothing saved. The transaction rolled back.",
        15, WHITE, False, False, SANS)],
      [("Every response shows the right status code (200 / 201 / 400 / 404) — that's the API contract working.",
        13, MUTED, False, True, SANS)]],
     space_after=4)

# ============================================================ SLIDE 4 — CODE SPOTLIGHT
s = prs.slides.add_slide(BLANK); bg(s, NAVY)
text(s, Inches(0.9), Inches(0.45), Inches(11.5), Inches(0.8),
     [[("💡  Code Spotlight — the atomic transaction", 30, WHITE, True, False, SANS)]])
cb = box(s, Inches(0.9), Inches(1.4), Inches(7.4), Inches(5.3), fill=CODEBG, radius=True)
code = [
    ("return await prisma.$transaction(async (tx) => {", CODEFG),
    ("  const order = await tx.order.create({", CODEFG),
    ("    data: { customer, status, totalPrice }", CODEFG),
    ("  })", CODEFG),
    ("  await tx.orderItem.createMany({", CODEFG),
    ("    data: items.map(d =>", CODEFG),
    ("      ({ ...d, orderId: order.id }))", CODEFG),
    ("  })", CODEFG),
    ("  return tx.order.findUnique({", CODEFG),
    ("    where: { id: order.id },", CODEFG),
    ("    include: { orderItems:", CODEFG),
    ("      { include: { product: true } } }", CODEFG),
    ("  })", CODEFG),
    ("})  // any error → the WHOLE thing rolls back", GREENC),
]
tf = cb.text_frame; tf.word_wrap=True; tf.margin_left=Inches(0.25); tf.margin_top=Inches(0.2)
for i,(ln,col) in enumerate(code):
    p = tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.line_spacing=1.05
    r=p.add_run(); r.text=ln if ln else " "
    r.font.size=Pt(14); r.font.name=MONO; r.font.color.rgb=col
    if "rolls back" in ln: r.font.italic=True
box(s, Inches(8.5), Inches(1.4), Inches(3.95), Inches(5.3), fill=INK, line=ACCENT, line_w=Pt(1.5), radius=True)
text(s, Inches(8.75), Inches(1.6), Inches(3.5), Inches(5.0),
     [[("Why it matters", 18, ACCENT, True, False, SANS)],
      [("An order = several writes:", 14, WHITE, False, False, SANS)],
      [("the order + a row per item.", 14, WHITE, False, False, SANS)],
      [("", 6, WHITE, False, False, SANS)],
      [("If item 3 fails → a half-made", 14, WHITE, False, False, SANS)],
      [("order + a wrongly-charged", 14, WHITE, False, False, SANS)],
      [("customer.", 14, WHITE, False, False, SANS)],
      [("", 6, WHITE, False, False, SANS)],
      [("$transaction = all-or-nothing.", 14, GREENC, True, False, SANS)],
      [("", 6, WHITE, False, False, SANS)],
      [("Prices come from the DB,", 14, WHITE, False, False, SANS)],
      [("never the client — no buying", 14, WHITE, False, False, SANS)],
      [("a laptop for one cent.", 14, YELLOW, False, True, SANS)]],
     space_after=4)

# ============================================================ SLIDE 5 — REFLECTION
s = prs.slides.add_slide(BLANK); bg(s, NAVY)
text(s, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.9),
     [[("🔭  Reflection", 40, WHITE, True, False, SANS)]])
cards = [
    ("⭐ Favorite", ACCENT,
     "The atomic POST /orders transaction — the one place all three tables meet. Making it "
     "bulletproof (rollback on failure, price server-side) taught me real backend safety."),
    ("🧗 Most Challenging", GREEN,
     "Relationships & cascade deletes. Foreign keys confused me — the raw orderId column vs. the "
     "order relation field, and that Prisma needs BOTH ends of every relation. Then it clicked."),
    ("🚀 Next Steps", YELLOW,
     "Deploy to Render for a live URL · build a Past Orders page + filter by email · add real "
     "customer accounts so customer isn't a placeholder."),
]
cw = Inches(3.7); cx = Inches(0.9); cy = Inches(1.9); cgap = Inches(0.25)
for i,(title,col,body) in enumerate(cards):
    bx = Emu(int(cx)+i*(int(cw)+int(cgap)))
    box(s, bx, cy, cw, Inches(4.3), fill=INK, line=col, line_w=Pt(2), radius=True)
    box(s, bx, cy, cw, Inches(0.12), fill=col)
    text(s, Emu(int(bx)+Emu(220000)), Emu(int(cy)+Emu(300000)), Inches(3.3), Inches(0.7),
         [[(title, 20, col, True, False, SANS)]])
    text(s, Emu(int(bx)+Emu(220000)), Emu(int(cy)+Emu(1050000)), Inches(3.3), Inches(3.0),
         [[(body, 15, WHITE, False, False, SANS)]], line_spacing=1.1)

out = __file__.rsplit("/",1)[0] + "/StudentStore-Demo.pptx"
prs.save(out)
print("saved:", out, "| slides:", len(prs.slides._sldIdLst))
